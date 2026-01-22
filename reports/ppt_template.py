from pptx import Presentation

def create_ppt(report_data, output_file):
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Media Plan Recommendation"
    slide.placeholders[1].text = report_data["client"]

    # Recommendation Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommended Media Strategy"
    content = slide.placeholders[1]

    content.text = (
        f"Platform: {report_data['platform']}\n"
        f"Time Band: {report_data['time_band']}\n"
        f"Region: {report_data['region']}\n"
        f"Expected ROI: {report_data['expected_roi']}\n"
        f"Confidence: {report_data['confidence']}"
    )

    # Explanation Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Why This Works"
    slide.placeholders[1].text = report_data["explanation"]

    prs.save(output_file)
